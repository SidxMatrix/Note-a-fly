import os
import json
import re
import time
import traceback
import tempfile

import requests
import httpx
from dotenv import load_dotenv
from fastapi import FastAPI, HTTPException, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import Optional

from supabase import create_client

load_dotenv()

app = FastAPI(title="NoteFlow AI Service")

app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:3000",
        os.getenv("FRONTEND_URL", ""),
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ============ CONFIG ============
OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY", "")
OPENROUTER_URL = "https://openrouter.ai/api/v1/chat/completions"

SUPABASE_URL = os.getenv("FASTAPI_SUPABASE_URL", "")
SUPABASE_KEY = os.getenv("FASTAPI_SUPABASE_SERVICE_KEY", "")

supabase = create_client(SUPABASE_URL, SUPABASE_KEY)

# FFmpeg explicit path — avoids PATH dependency for yt-dlp,
# AND injected into os.environ["PATH"] so Whisper can find ffmpeg too.
FFMPEG_PATH = r"C:\Users\ASUS\AppData\Local\Microsoft\WinGet\Packages\Gyan.FFmpeg_Microsoft.Winget.Source_8wekyb3d8bbwe\ffmpeg-8.1-full_build\bin"

# Inject FFmpeg into PATH at startup so Whisper (and any subprocess) can find it.
# This must happen before any Whisper or audio calls.
if FFMPEG_PATH and FFMPEG_PATH not in os.environ.get("PATH", ""):
    os.environ["PATH"] = FFMPEG_PATH + os.pathsep + os.environ.get("PATH", "")
    print(f"✅ FFmpeg injected into PATH: {FFMPEG_PATH}")

# Content limits
MAX_TOTAL_CHARS = 25000
MAX_CHARS_PER_FILE = 25000
MIN_CHARS_PER_FILE = 5000

# Model cache
_free_models_cache: list[str] = []
_free_models_fetched_at: float = 0
FREE_MODELS_CACHE_TTL = 600

# Keywords to avoid — these models can't do instruction following
AVOID_KEYWORDS = [
    "coder", "code", "vision", "image", "embed",
    "moderation", "tts", "dall", "thinking", "reasoner",
    "reasoning", "r1", "r2", "r3", "qwq", "deepseek-r",
    "o1", "o3", "nemotron", "reflection",
]

# Preferred models tried first if available
PREFERRED_MODELS = [
    "meta-llama/llama-3.3-70b-instruct:free",
    "google/gemma-3-27b-it:free",
    "google/gemma-3-12b-it:free",      
    "minimax/minimax-m2.5:free",       
    "meta-llama/llama-3.1-70b-instruct:free",
    "meta-llama/llama-3.2-3b-instruct:free",
    "google/gemma-3-4b-it:free",
    "mistralai/mistral-7b-instruct:free",
]

# Whisper model — lazy loaded on first use
_whisper_model = None


# ============ WHISPER ============
def get_whisper_model():
    """Lazy-load Whisper model — only loads when first needed."""
    global _whisper_model
    if _whisper_model is None:
        import whisper
        print("⏳ Loading Whisper model (base)...")
        _whisper_model = whisper.load_model("base")
        print("✅ Whisper model loaded!")
    return _whisper_model


# ============ PYDANTIC MODELS ============
class ProcessRequest(BaseModel):
    session_id: str
    source_type: str          # "text" | "audio" | "file" | "youtube"
    file_url: Optional[str] = None
    content: Optional[str] = None
    user_id: str


class ChatRequest(BaseModel):
    message: str
    history: list[dict] = []


# ============ YOUTUBE UTILITIES ============
def is_youtube_url(url: str) -> bool:
    """Check if a string is a valid YouTube URL."""
    patterns = [
        r'(https?://)?(www\.)?youtube\.com/watch\?v=[\w-]+',
        r'(https?://)?(www\.)?youtu\.be/[\w-]+',
        r'(https?://)?(www\.)?youtube\.com/embed/[\w-]+',
        r'(https?://)?(www\.)?youtube\.com/v/[\w-]+',
    ]
    return any(re.match(p, url.strip()) for p in patterns)


def extract_video_id(url: str) -> Optional[str]:
    """Extract video ID from any YouTube URL format."""
    patterns = [
        r'(?:youtube\.com/watch\?v=|youtu\.be/|youtube\.com/embed/|youtube\.com/v/)([a-zA-Z0-9_-]{11})',
    ]
    for pattern in patterns:
        match = re.search(pattern, url)
        if match:
            return match.group(1)
    return None


def get_video_info(video_id: str) -> Optional[dict]:
    """Get video metadata using yt-dlp."""
    try:
        import yt_dlp
        ydl_opts = {
            "quiet": True,
            "no_warnings": True,
            "ffmpeg_location": FFMPEG_PATH,
        }
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(
                f"https://www.youtube.com/watch?v={video_id}",
                download=False,
            )
            return {
                "title": info.get("title", "Unknown Title"),
                "duration": info.get("duration", 0),
                "channel": info.get("channel", "Unknown Channel"),
            }
    except Exception as e:
        print(f"⚠️ Error getting video info: {e}")
        return None


def get_transcript(video_id: str) -> Optional[str]:
    """
    Fetch transcript via youtube-transcript-api >= 0.6.0.

    The new API removed the class-method interface. The correct approach is:
        api = YouTubeTranscriptApi()
        transcript_list = api.list(video_id)       # list available transcripts
        transcript = transcript_list.find_transcript(['en'])
        fetched = transcript.fetch()               # returns FetchedTranscript

    FetchedTranscript is iterable; each element has a .text attribute (not dict key).
    We also handle translation for non-English videos.
    """
    try:
        from youtube_transcript_api import YouTubeTranscriptApi, NoTranscriptFound, TranscriptsDisabled

        api = YouTubeTranscriptApi()

        # ── Attempt 1: find an English transcript directly ───────────────────
        try:
            transcript_list = api.list(video_id)
            transcript = transcript_list.find_transcript(["en", "en-US", "en-GB"])
            fetched = transcript.fetch()
            # FetchedTranscript elements expose .text (not dict["text"])
            full_text = " ".join(
                snippet.text for snippet in fetched
            ).strip()
            if full_text:
                print(f"✅ English transcript found ({len(full_text)} chars)")
                return full_text
        except NoTranscriptFound:
            print("⚠️ No English transcript — trying auto-generated...")
        except Exception as e1:
            print(f"⚠️ English transcript attempt failed: {e1}")

        # ── Attempt 2: find auto-generated English transcript ────────────────
        try:
            transcript_list = api.list(video_id)
            # find_generated_transcript raises NoTranscriptFound if none exist
            transcript = transcript_list.find_generated_transcript(["en", "en-US", "en-GB"])
            fetched = transcript.fetch()
            full_text = " ".join(snippet.text for snippet in fetched).strip()
            if full_text:
                print(f"✅ Auto-generated transcript found ({len(full_text)} chars)")
                return full_text
        except NoTranscriptFound:
            print("⚠️ No auto-generated English transcript — trying any language + translate...")
        except Exception as e2:
            print(f"⚠️ Auto-generated transcript attempt failed: {e2}")

        # ── Attempt 3: take whatever's available and translate to English ────
        try:
            transcript_list = api.list(video_id)
            # Iterate over all available transcripts and pick the first one
            transcript = None
            for t in transcript_list:
                transcript = t
                break

            if transcript is None:
                print("⚠️ No transcripts at all for this video")
                return None

            # Translate if not already English
            if transcript.language_code not in ("en", "en-US", "en-GB"):
                print(f"⚠️ Translating from {transcript.language_code} to English...")
                transcript = transcript.translate("en")

            fetched = transcript.fetch()
            full_text = " ".join(snippet.text for snippet in fetched).strip()
            if full_text:
                print(f"✅ Translated transcript found ({len(full_text)} chars)")
                return full_text

        except TranscriptsDisabled:
            print("⚠️ Transcripts are disabled for this video")
            return None
        except Exception as e3:
            print(f"⚠️ Translate/any-language transcript attempt failed: {e3}")

        return None

    except Exception as e:
        print(f"⚠️ Transcript error: {e}")
        traceback.print_exc()
        return None


def download_and_transcribe_youtube_audio(video_id: str) -> Optional[str]:
    """
    Download YouTube audio and transcribe with local Whisper.
    Used only when no transcript is available and video is under 25 minutes.

    FFmpeg is already in PATH (injected at startup), so Whisper will find it.
    yt-dlp also receives ffmpeg_location explicitly as a belt-and-suspenders measure.
    """
    import yt_dlp
    temp_base = tempfile.mktemp()
    try:
        ydl_opts = {
            "format": "bestaudio/best",
            "outtmpl": temp_base,
            "postprocessors": [{
                "key": "FFmpegExtractAudio",
                "preferredcodec": "mp3",
                "preferredquality": "128",
            }],
            "quiet": True,
            "no_warnings": True,
            "ffmpeg_location": FFMPEG_PATH,
        }

        print(f"📥 Downloading audio for: {video_id}")
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            ydl.download([f"https://www.youtube.com/watch?v={video_id}"])

        # Find the downloaded file — yt-dlp may append extension
        actual_path = None
        for ext in [".mp3", ".m4a", ".webm", ".opus", ""]:
            candidate = temp_base + ext
            if os.path.exists(candidate):
                actual_path = candidate
                break

        if not actual_path:
            print("❌ Audio file not found after download")
            return None

        file_size = os.path.getsize(actual_path)
        print(f"🎙️ Transcribing YouTube audio ({file_size} bytes)...")
        start = time.time()
        model = get_whisper_model()
        # fp16=False for CPU/CUDA compatibility; Whisper will use its own ffmpeg via PATH
        result = model.transcribe(actual_path, fp16=False)
        print(f"✅ Transcription done in {round(time.time() - start, 2)}s")
        return result.get("text", "").strip() or None

    except Exception as e:
        print(f"❌ YouTube audio transcription error: {e}")
        traceback.print_exc()
        return None
    finally:
        for ext in [".mp3", ".m4a", ".webm", ".opus", ""]:
            try:
                path = temp_base + ext
                if os.path.exists(path):
                    os.unlink(path)
            except Exception:
                pass


def process_youtube_url(url: str) -> tuple[Optional[str], Optional[str]]:
    """
    Full YouTube pipeline.
    Returns (content, error_message). Always returns a tuple — never raises.
    """
    try:
        if not is_youtube_url(url):
            return None, "Invalid YouTube URL"

        video_id = extract_video_id(url)
        if not video_id:
            return None, "Could not extract video ID"

        print(f"🎥 Processing YouTube video: {video_id}")

        video_info = get_video_info(video_id)
        if not video_info:
            return None, "Could not fetch video information"

        duration = video_info.get("duration", 0)
        title    = video_info.get("title", "Unknown")
        print(f"📺 Video: {title} ({duration}s)")

        # Try transcript first — always free, fast, no audio download needed
        transcript = get_transcript(video_id)
        if transcript:
            return f"[YouTube: {title}]\n\n{transcript}", None

        # No transcript — only attempt Whisper if under 25 minutes
        print("⚠️ No transcript available — attempting Whisper transcription...")
        if duration > 1500:
            return None, (
                f"Video '{title}' has no transcript and is too long to transcribe "
                f"({duration // 60} min). Use a video under 25 minutes."
            )

        content = download_and_transcribe_youtube_audio(video_id)
        if not content:
            return None, (
                "Could not transcribe video audio. "
                "Make sure FFmpeg is installed and accessible."
            )

        return f"[YouTube: {title}]\n\n{content}", None

    except Exception as e:
        print(f"❌ process_youtube_url error: {e}")
        traceback.print_exc()
        return None, f"YouTube processing failed: {str(e)}"


# ============ SMART TRUNCATION ============
def smart_truncate_text(text: str, max_chars: int) -> tuple[str, bool]:
    """
    Truncate text at a clean paragraph or sentence boundary where possible.
    Returns (truncated_text, was_truncated).
    """
    if len(text) <= max_chars:
        return text, False

    truncated = text[:max_chars]

    last_para = truncated.rfind("\n\n")
    if last_para > max_chars * 0.7:
        return truncated[:last_para] + "\n\n[... content truncated ...]", True

    last_sentence = max(truncated.rfind(". "), truncated.rfind(".\n"))
    if last_sentence > max_chars * 0.8:
        return truncated[:last_sentence + 1] + "\n\n[... content truncated ...]", True

    return truncated + "\n\n[... content truncated ...]", True


def combine_multiple_files(
    file_texts: list[str],
    file_names: list[str],
    max_total_chars: int = MAX_TOTAL_CHARS,
) -> str:
    """
    Combine multiple file contents with proportional allocation so every file
    is represented in the output even when total exceeds the character limit.
    """
    num_files = len(file_texts)
    if num_files == 0:
        return ""

    if num_files == 1:
        text, _ = smart_truncate_text(file_texts[0], max_total_chars)
        return f"--- FILE: {file_names[0]} ---\n\n{text}"

    total_original = sum(len(t) for t in file_texts)
    overhead_per_file = 150
    available_chars = max_total_chars - (overhead_per_file * num_files)

    chars_per_file = []
    for text in file_texts:
        proportion = len(text) / total_original if total_original > 0 else 1 / num_files
        allocated = int(available_chars * proportion)
        allocated = max(MIN_CHARS_PER_FILE, min(MAX_CHARS_PER_FILE, allocated))
        chars_per_file.append(allocated)

    while sum(chars_per_file) > available_chars and len(chars_per_file) > 0:
        max_idx = chars_per_file.index(max(chars_per_file))
        chars_per_file[max_idx] = max(MIN_CHARS_PER_FILE, chars_per_file[max_idx] - 1000)

    parts = [f"COMBINED CONTENT FROM {num_files} FILES — COVER ALL SECTIONS IN OUTPUT\n"]

    for i, (text, filename, max_chars) in enumerate(
        zip(file_texts, file_names, chars_per_file), 1
    ):
        truncated, was_truncated = smart_truncate_text(text, max_chars)
        parts.append(f"\n{'='*60}")
        parts.append(f"FILE {i} OF {num_files}: {filename}")
        parts.append(f"{'='*60}\n")
        parts.append(truncated)
        if was_truncated:
            print(f"⚠️ '{filename}' truncated: {len(text):,} → {len(truncated):,} chars")
        else:
            print(f"✅ '{filename}' included fully: {len(text):,} chars")

    parts.append(f"\n{'='*60}")
    parts.append(f"END OF ALL {num_files} FILES")
    parts.append(f"{'='*60}")

    return "\n".join(parts)


# ============ AI CALL ============
def fetch_free_models() -> list[str]:
    """
    Fetch all available free non-reasoning models from OpenRouter.
    Scores by context length and instruction-following indicators.
    Results cached for 10 minutes.
    """
    global _free_models_cache, _free_models_fetched_at

    now = time.time()
    if _free_models_cache and (now - _free_models_fetched_at) < FREE_MODELS_CACHE_TTL:
        return _free_models_cache

    print("🔍 Fetching available free models from OpenRouter...")
    try:
        resp = requests.get(
            "https://openrouter.ai/api/v1/models",
            headers={"Authorization": f"Bearer {OPENROUTER_API_KEY}"},
            timeout=15,
        )
        resp.raise_for_status()
        all_models = resp.json().get("data", [])

        scored: list[dict] = []
        for m in all_models:
            mid = m.get("id", "")
            if not mid.endswith(":free"):
                continue
            if not m.get("context_length"):
                continue

            mid_lower = mid.lower()
            if any(kw in mid_lower for kw in AVOID_KEYWORDS):
                continue

            score = min(m.get("context_length", 0) // 10000, 20)
            if "instruct" in mid_lower:
                score += 5

            scored.append({"id": mid, "score": score})

        scored.sort(key=lambda x: x["score"], reverse=True)
        all_ids = [s["id"] for s in scored]

        def sort_key(mid: str) -> tuple:
            try:
                return (0, PREFERRED_MODELS.index(mid))
            except ValueError:
                return (1, -next((s["score"] for s in scored if s["id"] == mid), 0))

        all_ids.sort(key=sort_key)

        print(f"✅ Found {len(all_ids)} free non-reasoning models")
        _free_models_cache = all_ids
        _free_models_fetched_at = now
        return all_ids

    except Exception as e:
        print(f"⚠️ Could not fetch models: {e}")
        return PREFERRED_MODELS


def call_openrouter(system_prompt: str, user_message: str, max_tokens: int = 4000) -> str:
    """
    Call OpenRouter, dynamically fetching available free models and retrying
    down the list until one succeeds. Strips leaked reasoning tags.
    """
    headers = {
        "Authorization": f"Bearer {OPENROUTER_API_KEY}",
        "Content-Type": "application/json",
        "HTTP-Referer": "http://localhost:3000",
        "X-Title": "Note-a-fly",
    }

    models = fetch_free_models()
    if not models:
        raise Exception("No free models available")

    for i, model_id in enumerate(models):
        print(f"   Trying model: {model_id}")
        try:
            if i > 0:
                time.sleep(2.5)

            payload = {
                "model": model_id,
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_message},
                ],
                "max_tokens": max_tokens,
                "temperature": 0.3,
            }

            resp = requests.post(
                OPENROUTER_URL,
                headers=headers,
                json=payload,
                timeout=120,
            )

            if resp.status_code == 429:
                print(f"  ⚠️ Rate limited, trying next...")
                continue

            if resp.status_code in (404, 400):
                print(f"  ⚠️ HTTP {resp.status_code}, trying next...")
                continue

            if resp.status_code == 200:
                body = resp.json()
                choices = body.get("choices", [])
                if choices:
                    content = choices[0].get("message", {}).get("content", "")
                    if content and "<think>" in content:
                        if "</think>" in content:
                            content = content.split("</think>")[-1].strip()
                        else:
                            content = content.split("<think>")[0].strip()
                    if content and content.strip():
                        actual_model = body.get("model", model_id)
                        print(f"  ✅ Success with: {actual_model}")
                        return content.strip()

            print(f"  ⚠️ Failed: HTTP {resp.status_code}")

        except Exception as e:
            print(f"  ⚠️ Error with {model_id}: {str(e)}")
            continue

    raise Exception("All available free models failed or are rate limited")


# ============ PROMPTS ============
NOTES_PROMPT = """You are an expert academic note-taking assistant.
Transform raw content into clean, well-organized study notes.

Rules:
- Use markdown headings (##, ###) for sections
- Use bullet points for key ideas
- **Bold** important terms, definitions, concepts
- Keep sentences short and scannable
- If input contains multiple files or sections, create notes covering ALL of them
- Add "💡 Key Insight:" for crucial concepts
- Add "⚠️ Exam Alert:" for likely exam topics
- End with "## 📝 Key Takeaways" section (5-7 most important points)
- Be comprehensive but concise"""

# ============ PROMPTS ============
FLASHCARDS_PROMPT = """You are an AI study assistant that creates flashcards.

Generate 8-12 flashcards from the given content.

STRICT FORMAT - return ONLY valid JSON:
{
  "flashcards": [
    {"front": "question here", "back": "concise answer here", "difficulty": "easy"},
    {"front": "question here", "back": "concise answer here", "difficulty": "medium"},
    {"front": "question here", "back": "concise answer here", "difficulty": "hard"}
  ]
}

Rules:
- Each flashcard tests ONE concept
- If multiple files exist, ensure flashcards come from ALL of them
- Questions should require understanding, not just memorization
- Answers MUST be 1-2 sentences MAX — absolutely no longer. Be brutal about brevity.
- Front (question) MUST be under 12 words
- Back (answer) MUST be under 25 words
- Difficulty: easy (recall), medium (understand), hard (apply)
- Return ONLY the JSON object, nothing else"""

QUIZ_PROMPT = """You are an AI exam generator.

Generate a quiz with exactly 7 questions.

STRICT FORMAT - return ONLY valid JSON:
{
  "title": "Quiz: Topic Name",
  "questions": [
    {
      "id": "q1",
      "type": "mcq",
      "question": "question text",
      "options": ["A) option", "B) option", "C) option", "D) option"],
      "correct_answer": "B) option",
      "explanation": "why this is correct"
    },
    {
      "id": "q2",
      "type": "true_false",
      "question": "statement to evaluate",
      "options": ["True", "False"],
      "correct_answer": "False",
      "explanation": "why"
    },
    {
      "id": "q3",
      "type": "short_answer",
      "question": "question text",
      "options": [],
      "correct_answer": "expected answer",
      "explanation": "key points"
    }
  ]
}

Rules:
- Mix: 4 MCQs, 2 True/False, 1 Short Answer
- MCQs must have exactly 4 options
- Every question needs an explanation
- Return ONLY the JSON object"""

MINDMAP_PROMPT = """Convert this content into a mind map structure.

Return ONLY valid JSON:
{
  "title": "Main Topic",
  "nodes": [
    {"id": "1", "type": "central", "position": {"x": 400, "y": 300}, "data": {"label": "Main Topic"}},
    {"id": "2", "type": "branch", "position": {"x": 150, "y": 150}, "data": {"label": "Subtopic 1"}},
    {"id": "3", "type": "branch", "position": {"x": 650, "y": 150}, "data": {"label": "Subtopic 2"}},
    {"id": "4", "type": "leaf", "position": {"x": 50, "y": 50}, "data": {"label": "Detail 1"}},
    {"id": "5", "type": "leaf", "position": {"x": 250, "y": 50}, "data": {"label": "Detail 2"}}
  ],
  "edges": [
    {"id": "e1-2", "source": "1", "target": "2"},
    {"id": "e1-3", "source": "1", "target": "3"},
    {"id": "e2-4", "source": "2", "target": "4"},
    {"id": "e2-5", "source": "2", "target": "5"}
  ]
}

Rules:
- 1 central node, 3-5 branches, 2-3 leaves per branch
- Keep labels SHORT (2-5 words)
- Spread positions in a radial layout
- Return ONLY JSON"""

IMPORTANT_QUESTIONS_PROMPT = """You are a professor creating exam questions.

Generate 6-8 important exam questions from this content.

Return ONLY valid JSON:
{
  "questions": [
    {
      "question": "question text",
      "expected_answer": "model answer",
      "type": "short",
      "estimated_marks": 3,
      "topic": "topic name"
    },
    {
      "question": "question text",
      "expected_answer": "detailed model answer",
      "type": "long",
      "estimated_marks": 8,
      "topic": "topic name"
    }
  ]
}

Rules:
- Mix: 3 short (2-3 marks), 3 long (5-10 marks), 1-2 medium
- Include expected answers
- Return ONLY JSON"""


# ============ JSON PARSER ============
def parse_json_response(text: str) -> dict:
    """Extract JSON from AI response, handling markdown code blocks."""
    text = text.strip()
    if text.startswith("```json"):
        text = text[7:]
    elif text.startswith("```"):
        text = text[3:]
    if text.endswith("```"):
        text = text[:-3]
    text = text.strip()

    try:
        return json.loads(text)
    except json.JSONDecodeError:
        start = text.find("{")
        end = text.rfind("}") + 1
        if start >= 0 and end > start:
            try:
                return json.loads(text[start:end])
            except json.JSONDecodeError:
                pass
        return {}


# ============ FILE PARSING ============
def parse_pdf(file_path: str) -> str:
    import pdfplumber
    text_parts = []
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages):
            page_text = page.extract_text()
            if page_text:
                text_parts.append(f"--- Page {i + 1} ---\n{page_text}")
    return "\n\n".join(text_parts)


def parse_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    text_parts = []
    for i, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            text_parts.append(f"--- Slide {i + 1} ---\n" + "\n".join(slide_texts))
    return "\n\n".join(text_parts)


def parse_docx(file_path: str) -> str:
    """Extract text from DOCX including table cells."""
    from docx import Document
    doc = Document(file_path)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                parts.append(row_text)
    return "\n\n".join(parts)


def parse_txt(file_path: str) -> str:
    """Extract text from plain text files with encoding fallback."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read().strip()
    except UnicodeDecodeError:
        with open(file_path, "r", encoding="latin-1") as f:
            return f.read().strip()


async def download_and_parse_file(file_url: str, file_type: str) -> str:
    """Download a file from Supabase storage and extract its text."""
    async with httpx.AsyncClient() as client:
        response = await client.get(file_url, timeout=60)
        response.raise_for_status()

    suffix_map = {
        "pdf":   ".pdf",
        "pptx":  ".pptx",
        "docx":  ".docx",
        "txt":   ".txt",
        "image": ".jpg",
    }
    suffix = suffix_map.get(file_type, ".bin")

    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(response.content)
        tmp_path = tmp.name

    try:
        if file_type == "pdf":
            return parse_pdf(tmp_path)
        elif file_type == "pptx":
            return parse_pptx(tmp_path)
        elif file_type == "docx":
            return parse_docx(tmp_path)
        elif file_type == "txt":
            return parse_txt(tmp_path)
        else:
            return "Unsupported file type"
    finally:
        os.unlink(tmp_path)


# ============ ROUTES ============
@app.get("/health")
def health_check():
    return {"status": "healthy", "service": "noteflow-ai"}


@app.get("/api/models")
def list_models():
    """Debug — see what free models are currently available."""
    models = fetch_free_models()
    return {"count": len(models), "models": models}


@app.post("/api/transcribe")
async def transcribe_audio(audio: UploadFile = File(...)):
    """
    Transcribe an uploaded audio file using local Whisper.
    Accepts any audio format Whisper supports (webm, mp3, wav, m4a, ogg).
    FFmpeg is available via PATH injection at startup.
    """
    suffix = os.path.splitext(audio.filename)[1] if audio.filename else ".webm"

    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        content = await audio.read()
        tmp.write(content)
        tmp_path = tmp.name

    try:
        file_size = os.path.getsize(tmp_path)
        print(f"🎙️ Transcribing audio ({file_size} bytes)...")
        start = time.time()

        model = get_whisper_model()
        result = model.transcribe(tmp_path, fp16=False)
        elapsed = round(time.time() - start, 2)
        print(f"✅ Transcription done in {elapsed}s")

        text = result.get("text", "").strip()
        if not text:
            raise HTTPException(status_code=400, detail="No speech detected in audio")

        return {"text": text, "duration_seconds": elapsed}

    except HTTPException:
        raise
    except Exception as e:
        print(f"❌ Transcription error: {e}")
        raise HTTPException(status_code=500, detail=f"Transcription failed: {str(e)}")
    finally:
        os.unlink(tmp_path)


@app.post("/api/chat")
async def chat(req: ChatRequest):
    if not req.message or len(req.message.strip()) == 0:
        raise HTTPException(status_code=400, detail="Message required")

    if len(req.message.strip()) > 500:
        raise HTTPException(status_code=400, detail="Message too long")

    system_prompt = """You are the Note-a-fly assistant — an intelligent, warm, and concise guide for the Note-a-fly study platform ("The Living Manuscript").

You help users understand:
- What Note-a-fly is: an AI-powered study app that transforms lectures, PDFs, and text into structured notes, flashcards, exam questions, and mindmaps
- How to use it: paste text, upload files (PDF/PPTX/DOCX/TXT), record audio, or paste a YouTube URL — AI processes and generates a full study kit
- Features: ambient lecture recording, YouTube transcript extraction, multimodal file upload, auto-generated flashcards, quizzes, exam Q&A, mindmaps
- Pricing: free 14-day trial, no credit card required
- Tech: powered by state-of-the-art AI models via OpenRouter, local Whisper for audio, Supabase for secure storage
- Account/auth: sign up with email, sessions saved to personal library

Tone: Sophisticated, encouraging, precise. Never robotic. Max 3 sentences per response unless a list is genuinely needed. Do not answer questions unrelated to Note-a-fly or studying. If asked something outside scope, gently redirect."""

    history_text = ""
    for msg in req.history[-6:]:
        role    = msg.get("role", "")
        content = msg.get("content", "")
        if role and content:
            history_text += f"\n{role.upper()}: {content}"

    user_message = (
        f"{history_text}\nUSER: {req.message.strip()}"
        if history_text
        else req.message.strip()
    )

    try:
        reply = call_openrouter(system_prompt, user_message, max_tokens=300)
        return {"reply": reply}
    except Exception as e:
        print(f"Chat error: {e}")
        raise HTTPException(status_code=502, detail="AI service unavailable")


@app.post("/api/process")
async def process_session(request: ProcessRequest):
    try:
        print(f"\n{'='*50}")
        print(f" Processing session: {request.session_id}")
        print(f"   Type: {request.source_type}")
        print(f"{'='*50}")

        # ── STEP 1: Extract content ──────────────────────────────────────────
        content = ""

        if request.source_type == "text":
            content = request.content or ""

        elif request.source_type == "youtube":
            url = request.content or ""
            if not url:
                raise HTTPException(400, "content field must contain the YouTube URL")
            extracted, error = process_youtube_url(url)
            if error:
                supabase.table("sessions").update({"status": "failed"}).eq("id", request.session_id).execute()
                raise HTTPException(400, error)
            content = extracted or ""

        elif request.source_type in ("pdf", "pptx", "docx", "txt"):
            if not request.file_url:
                raise HTTPException(400, "file_url required")
            print(" Parsing file...")
            content = await download_and_parse_file(request.file_url, request.source_type)
            supabase.table("sessions").update(
                {"raw_transcript": content}
            ).eq("id", request.session_id).execute()

        elif request.source_type == "audio":
            content = request.content or ""
            if not content:
                supabase.table("sessions").update({"status": "failed"}).eq("id", request.session_id).execute()
                raise HTTPException(400, "Audio transcript content is empty")

        if not content or len(content.strip()) < 20:
            supabase.table("sessions").update({"status": "failed"}).eq("id", request.session_id).execute()
            raise HTTPException(400, "Content too short to process")

        content, was_truncated = smart_truncate_text(content, MAX_TOTAL_CHARS)
        if was_truncated:
            print(f"⚠️ Content truncated to {MAX_TOTAL_CHARS:,} chars")
        else:
            print(f" Content: {len(content):,} chars")

        # ── STEP 2: Notes ────────────────────────────────────────────────────
        print(" Generating notes...")
        notes_content = call_openrouter(
            NOTES_PROMPT,
            f"Create study notes from:\n\n{content}"
        )
        time.sleep(3)
        print(" Generating summary...")
        meta_text = call_openrouter(
            'Return ONLY a JSON object with: {"summary": "2-3 sentence summary", "key_topics": ["topic1", "topic2"], "exam_important_sections": [{"section": "name", "importance": "high", "reason": "why"}]}. Return ONLY JSON.',
            f"Analyze this and return JSON:\n\n{notes_content[:3000]}"
        )
        meta = parse_json_response(meta_text)

        supabase.table("notes").insert({
            "session_id": request.session_id,
            "user_id":    request.user_id,
            "content":    notes_content,
            "summary":    meta.get("summary", ""),
            "key_topics": meta.get("key_topics", []),
            "exam_important_sections": meta.get("exam_important_sections", []),
        }).execute()
        print("✅ Notes saved")
        time.sleep(3)
        # ── STEP 3: Flashcards ───────────────────────────────────────────────
        print("🃏 Generating flashcards...")
        flashcards_text = call_openrouter(
            FLASHCARDS_PROMPT,
            f"Create flashcards from:\n\n{notes_content[:5000]}"
        )
        flashcards_data = parse_json_response(flashcards_text)
        flashcards      = flashcards_data.get("flashcards", [])

        for card in flashcards:
            supabase.table("flashcards").insert({
                "session_id": request.session_id,
                "user_id":    request.user_id,
                "front":      card.get("front", ""),
                "back":       card.get("back", ""),
                "difficulty": card.get("difficulty", "medium"),
            }).execute()
        print(f"✅ {len(flashcards)} flashcards saved")
        time.sleep(3)
        # ── STEP 4: Quiz ─────────────────────────────────────────────────────
        print(" Generating quiz...")
        quiz_text = call_openrouter(
            QUIZ_PROMPT,
            f"Create a quiz from:\n\n{notes_content[:5000]}"
        )
        quiz_data = parse_json_response(quiz_text)
        questions = quiz_data.get("questions", [])

        if questions:
            supabase.table("quizzes").insert({
                "session_id":      request.session_id,
                "user_id":         request.user_id,
                "title":           quiz_data.get("title", "Quiz"),
                "questions":       questions,
                "total_questions": len(questions),
            }).execute()
        print(f"✅ Quiz saved ({len(questions)} questions)")
        time.sleep(3)
        # ── STEP 5: Mindmap ──────────────────────────────────────────────────
        print(" Generating mindmap...")
        mindmap_text = call_openrouter(
            MINDMAP_PROMPT,
            f"Create a mind map from:\n\n{notes_content[:4000]}"
        )
        mindmap_data = parse_json_response(mindmap_text)

        if mindmap_data.get("nodes"):
            supabase.table("mindmaps").insert({
                "session_id": request.session_id,
                "user_id":    request.user_id,
                "title":      mindmap_data.get("title", "Mind Map"),
                "nodes":      mindmap_data.get("nodes", []),
                "edges":      mindmap_data.get("edges", []),
            }).execute()
        print("✅ Mindmap saved")
        time.sleep(3)
        # ── STEP 6: Important questions ──────────────────────────────────────
        print(" Generating important questions...")
        questions_text = call_openrouter(
            IMPORTANT_QUESTIONS_PROMPT,
            f"Generate exam questions from:\n\n{notes_content[:5000]}"
        )
        imp_data = parse_json_response(questions_text)

        if imp_data.get("questions"):
            supabase.table("important_questions").insert({
                "session_id": request.session_id,
                "user_id":    request.user_id,
                "questions":  imp_data.get("questions", []),
            }).execute()
        print("✅ Important questions saved")
        time.sleep(3)
        # ── STEP 7: Mark complete ────────────────────────────────────────────
        supabase.table("sessions").update(
            {"status": "completed"}
        ).eq("id", request.session_id).execute()

        print(f"\n✅ Session {request.session_id} COMPLETE!")
        return {"status": "completed", "session_id": request.session_id}

    except HTTPException:
        raise
    except Exception as e:
        print(f"❌ Error: {e}")
        traceback.print_exc()
        try:
            supabase.table("sessions").update(
                {"status": "failed"}
            ).eq("id", request.session_id).execute()
        except Exception:
            pass
        raise HTTPException(500, str(e))


if __name__ == "__main__":
    import uvicorn
    print("\n NoteFlow AI Service starting...")
    print(f"   OpenRouter key: {'✅ Set' if OPENROUTER_API_KEY else '❌ Missing'}")
    print(f"   Supabase URL:   {'✅ Set' if SUPABASE_URL else '❌ Missing'}")
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)